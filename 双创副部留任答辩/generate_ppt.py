import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
import os

# === CONFIG ===
IMG_DIR = "图片素材"
OUTPUT = "孙杨_办公室副部长留任答辩.pptx"
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

# === COLOR PALETTE (Cyberpunk/Tech Dark Blue) ===
DARK_BG = RGBColor(0x0A, 0x0E, 0x1A)
DARK_BG2 = RGBColor(0x0F, 0x15, 0x2A)
DARK_CARD = RGBColor(0x14, 0x1E, 0x3A)
ACCENT_BLUE = RGBColor(0x00, 0xB4, 0xD8)      # Cyan accent
ACCENT_PURPLE = RGBColor(0x7B, 0x2F, 0xBE)     # Purple accent
ACCENT_PINK = RGBColor(0xE8, 0x3E, 0x8C)       # Pink accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xCC)
TEXT_GRAY = RGBColor(0x88, 0x8A, 0xA0)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GRADIENT_TOP = RGBColor(0x0F, 0x17, 0x2E)
GRADIENT_BOT = RGBColor(0x06, 0x0A, 0x14)

# === HELPER FUNCTIONS ===
def add_bg(slide, color=DARK_BG):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, corner_radius=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """Add a text box with single style."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, default_size=14, default_color=WHITE, default_bold=False, alignment=PP_ALIGN.LEFT, line_spacing=1.5):
    """Add a text box with multiple paragraphs, each can have different styles.
    lines: list of dicts or strings. dict keys: text, size, color, bold
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            line = {'text': line}
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line.get('text', '')
        p.font.size = Pt(line.get('size', default_size))
        p.font.color.rgb = line.get('color', default_color)
        p.font.bold = line.get('bold', default_bold)
        p.font.name = 'Microsoft YaHei'
        p.alignment = alignment
        p.space_after = Pt(line.get('space_after', 4))
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE, height=Pt(3)):
    """Add a horizontal accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_image_safe(slide, img_name, left, top, width=None, height=None):
    """Add image if exists."""
    path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(path):
        if width and height:
            return slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            return slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            return slide.shapes.add_picture(path, left, top, height=height)
        else:
            return slide.shapes.add_picture(path, left, top)
    return None

def add_section_number(slide, left, top, number, color=ACCENT_BLUE):
    """Add a large semi-transparent section number for visual effect."""
    txBox = slide.shapes.add_textbox(left, top, Inches(1.5), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(number).zfill(2)
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.LEFT
    # Make it slightly transparent via XML
    from pptx.oxml.ns import qn
    solidFill = txBox.text_frame.paragraphs[0].runs
    return txBox

def add_glass_card(slide, left, top, width, height, border_color=None):
    """Add a glassmorphism-style card."""
    card = add_rect(slide, left, top, width, height, fill_color=DARK_CARD)
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(0.5)
    else:
        card.line.fill.background()
    return card

# ======================================================================
# SLIDE 1: COVER / TITLE
# ======================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, DARK_BG)

# Background decorative elements
add_rect(slide1, Inches(0), Inches(0), Inches(13.333), Inches(0.08), fill_color=ACCENT_BLUE)
add_rect(slide1, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), fill_color=ACCENT_BLUE)

# Left decorative bar
add_rect(slide1, Inches(0.6), Inches(1.2), Inches(0.06), Inches(5.1), fill_color=ACCENT_BLUE)

# Big accent circle (decorative)
circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-0.8), Inches(5), Inches(5))
circle.fill.background()
circle.line.color.rgb = ACCENT_BLUE
circle.line.width = Pt(0.5)
circle2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(3.5), Inches(3), Inches(3))
circle2.fill.background()
circle2.line.color.rgb = ACCENT_PURPLE
circle2.line.width = Pt(0.5)

# Title text
add_textbox(slide1, Inches(1.2), Inches(1.8), Inches(8), Inches(0.8),
            "大学生创新创业俱乐部", font_size=20, color=LIGHT_GRAY, bold=False)

add_textbox(slide1, Inches(1.2), Inches(2.4), Inches(8), Inches(1.2),
            "换届竞选答辩", font_size=52, color=WHITE, bold=True)

add_accent_line(slide1, Inches(1.2), Inches(3.55), Inches(3.5), ACCENT_BLUE, Pt(4))

# Candidate info
add_textbox(slide1, Inches(1.2), Inches(4.0), Inches(6), Inches(0.6),
            "竞选人：孙  杨", font_size=26, color=WHITE, bold=True)
add_textbox(slide1, Inches(1.2), Inches(4.7), Inches(6), Inches(0.6),
            "竞选职位：办公室副部长", font_size=22, color=ACCENT_BLUE, bold=False)
add_textbox(slide1, Inches(1.2), Inches(5.4), Inches(6), Inches(0.5),
            "信息科学与技术学院 · 信息安全 2025-02 班", font_size=15, color=TEXT_GRAY)

# ======================================================================
# SLIDE 2: TABLE OF CONTENTS
# ======================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, DARK_BG)

add_textbox(slide2, Inches(1.2), Inches(0.8), Inches(4), Inches(0.7),
            "目  录", font_size=42, color=WHITE, bold=True)
add_textbox(slide2, Inches(1.2), Inches(1.4), Inches(4), Inches(0.4),
            "CONTENTS", font_size=16, color=ACCENT_BLUE, bold=False)
add_accent_line(slide2, Inches(1.2), Inches(1.9), Inches(2.5), ACCENT_BLUE, Pt(3))

sections = [
    ("01", "个人介绍", "基本信息、性格特点与获奖经历"),
    ("02", "我与双创", "在双创的成长故事与工作记录"),
    ("03", "工作能力", "经验积累、核心竞争力与留任原因"),
    ("04", "未来规划", "工作展望与社团发展计划"),
]

for i, (num, title, desc) in enumerate(sections):
    y_start = Inches(2.8) + Inches(i * 1.15)
    x_start = Inches(1.5)

    # Number
    add_textbox(slide2, x_start, y_start, Inches(0.8), Inches(0.8),
                num, font_size=40, color=ACCENT_BLUE, bold=True)

    # Vertical line
    add_rect(slide2, x_start + Inches(0.85), y_start + Inches(0.08), Inches(0.04), Inches(1.0), fill_color=RGBColor(0x1E, 0x30, 0x50))

    # Title + description
    add_textbox(slide2, x_start + Inches(1.1), y_start, Inches(6), Inches(0.5),
                title, font_size=22, color=WHITE, bold=True)
    add_textbox(slide2, x_start + Inches(1.1), y_start + Inches(0.45), Inches(7), Inches(0.4),
                desc, font_size=13, color=TEXT_GRAY)

# ======================================================================
# SLIDE 3: PERSONAL INTRODUCTION (个人介绍)
# ======================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, DARK_BG)

# Section header
add_textbox(slide3, Inches(1.0), Inches(0.5), Inches(4), Inches(0.6),
            "个人介绍", font_size=36, color=WHITE, bold=True)
add_textbox(slide3, Inches(1.0), Inches(1.0), Inches(4), Inches(0.3),
            "PERSONAL INTRODUCTION", font_size=13, color=ACCENT_BLUE)
add_accent_line(slide3, Inches(1.0), Inches(1.4), Inches(2.0), ACCENT_BLUE, Pt(3))

# Left side: Photo area
add_glass_card(slide3, Inches(1.0), Inches(2.0), Inches(4.5), Inches(5.0))

img_shape = add_image_safe(slide3, "自拍.jpg", Inches(1.3), Inches(2.3), width=Inches(3.9))
if not img_shape:
    # Placeholder
    placeholder = add_rect(slide3, Inches(1.3), Inches(2.3), Inches(3.9), Inches(4.5), fill_color=DARK_CARD)
    add_textbox(slide3, Inches(1.3), Inches(4.0), Inches(3.9), Inches(0.5),
                "[个人照片]", font_size=16, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Right side: Basic info cards
info_items = [
    ("姓  名", "孙  杨"),
    ("学  院", "信息科学与技术学院"),
    ("专  业", "信息安全 2025-02 班"),
    ("现任职务", "大学生创新创业俱乐部 办公室干事"),
    ("竞选职务", "办公室副部长"),
]

for i, (label, value) in enumerate(info_items):
    y = Inches(2.0) + Inches(i * 1.0)
    # Label
    add_textbox(slide3, Inches(6.2), y, Inches(1.3), Inches(0.4),
                label, font_size=14, color=ACCENT_BLUE, bold=True)
    # Value
    add_textbox(slide3, Inches(7.6), y, Inches(5), Inches(0.4),
                value, font_size=16, color=WHITE, bold=False)
    # Divider line
    add_rect(slide3, Inches(6.2), y + Inches(0.55), Inches(6.0), Inches(0.015), fill_color=RGBColor(0x1E, 0x30, 0x50))

# Other roles
add_textbox(slide3, Inches(6.2), Inches(7.0), Inches(5), Inches(0.3),
            "其他职务：班级学习委员 | 学风领航中心学业指导部成员", font_size=12, color=TEXT_GRAY)

# ======================================================================
# SLIDE 4: PERSONALITY & AWARDS (性格特点 & 获奖情况)
# ======================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, DARK_BG)

add_textbox(slide4, Inches(1.0), Inches(0.5), Inches(4), Inches(0.6),
            "性格特点 & 获奖经历", font_size=36, color=WHITE, bold=True)
add_textbox(slide4, Inches(1.0), Inches(1.0), Inches(4), Inches(0.3),
            "PERSONALITY & AWARDS", font_size=13, color=ACCENT_BLUE)
add_accent_line(slide4, Inches(1.0), Inches(1.4), Inches(2.0), ACCENT_BLUE, Pt(3))

# Personality traits as tags
traits = [
    ("🎯", "认真负责", "高效执行每项任务"),
    ("🤝", "善于沟通", "跨部门协调对接"),
    ("🔥", "积极开朗", "主动融入团队氛围"),
    ("💡", "敢于尝试", "拥抱挑战与变化"),
    ("📋", "统筹协调", "多任务并行管理"),
    ("🎨", "团队协作", "凝聚团队向心力"),
]

for i, (icon, title, desc) in enumerate(traits):
    col = i % 3
    row = i // 3
    x = Inches(1.0) + Inches(col * 4.0)
    y = Inches(1.9) + Inches(row * 1.8)

    card = add_glass_card(slide4, x, y, Inches(3.6), Inches(1.5))
    add_textbox(slide4, x + Inches(0.3), y + Inches(0.15), Inches(3.0), Inches(0.5),
                f"{icon}  {title}", font_size=18, color=WHITE, bold=True)
    add_textbox(slide4, x + Inches(0.3), y + Inches(0.75), Inches(3.0), Inches(0.4),
                desc, font_size=13, color=TEXT_GRAY)

# Awards section
awards_y = Inches(5.9)
add_glass_card(slide4, Inches(1.0), awards_y, Inches(11.3), Inches(1.2))

add_textbox(slide4, Inches(1.4), awards_y + Inches(0.1), Inches(2), Inches(0.4),
            "🏆 获奖经历", font_size=17, color=GOLD, bold=True)

add_multiline_textbox(slide4, Inches(1.4), awards_y + Inches(0.55), Inches(10.5), Inches(0.6), [
    {'text': '● 2026年5月  获评大学生创新创业俱乐部「优秀干事」称号', 'size': 15, 'color': WHITE, 'bold': False},
    {'text': '● 2026年4月  获集体荣誉：西南交通大学太极拳比赛二等奖', 'size': 15, 'color': WHITE, 'bold': False},
])

# ======================================================================
# SLIDE 5: ME & THE CLUB (我与双创)
# ======================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, DARK_BG)

add_textbox(slide5, Inches(1.0), Inches(0.5), Inches(4), Inches(0.6),
            "我与双创", font_size=36, color=WHITE, bold=True)
add_textbox(slide5, Inches(1.0), Inches(1.0), Inches(4), Inches(0.3),
            "MY JOURNEY WITH THE CLUB", font_size=13, color=ACCENT_BLUE)
add_accent_line(slide5, Inches(1.0), Inches(1.4), Inches(2.0), ACCENT_BLUE, Pt(3))

# Timeline style journey
journey = [
    ("初入双创", "怀揣对创新创业的热情加入双创大家庭，成为办公室干事"),
    ("融入成长", "参与奶茶会、全体大会，结识志同道合的伙伴，迅速融入团队"),
    ("担当历练", "负责会议记录、物资管理、人员协调，逐步成长为可靠执行者"),
    ("收获认可", "获评「优秀干事」，团队获太极拳比赛二等奖，见证个人与集体共同成长"),
]

for i, (title, desc) in enumerate(journey):
    y = Inches(2.0) + Inches(i * 1.25)
    x = Inches(1.5)

    # Timeline dot
    dot = slide5.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.1), y + Inches(0.12), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_BLUE if i < 3 else GOLD
    dot.line.fill.background()

    # Timeline line (except last)
    if i < 3:
        add_rect(slide5, x + Inches(0.17), y + Inches(0.35), Inches(0.03), Inches(0.95), fill_color=RGBColor(0x1E, 0x30, 0x50))

    # Content card
    add_glass_card(slide5, x + Inches(0.6), y, Inches(10.0), Inches(1.0))
    add_textbox(slide5, x + Inches(0.9), y + Inches(0.08), Inches(3), Inches(0.4),
                title, font_size=18, color=ACCENT_BLUE, bold=True)
    add_textbox(slide5, x + Inches(0.9), y + Inches(0.5), Inches(9.3), Inches(0.4),
                desc, font_size=13, color=TEXT_GRAY)

# ======================================================================
# SLIDE 6: WORK RECORDS (工作记录)
# ======================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, DARK_BG)

add_textbox(slide6, Inches(1.0), Inches(0.5), Inches(4), Inches(0.6),
            "工作记录", font_size=36, color=WHITE, bold=True)
add_textbox(slide6, Inches(1.0), Inches(1.0), Inches(4), Inches(0.3),
            "WORK RECORDS", font_size=13, color=ACCENT_BLUE)
add_accent_line(slide6, Inches(1.0), Inches(1.4), Inches(2.0), ACCENT_BLUE, Pt(3))

work_items = [
    ("📝", "日常会议记录", "负责办公室例会及大型活动的会议纪要整理与归档，确保信息传达准确无遗漏"),
    ("📦", "活动物资管理", "统筹大型双创活动的物资采购、调配与清点，保障活动物资链高效运转"),
    ("👥", "人员协调调度", "协助进行活动现场的人员排班与调度，确保各环节衔接顺畅"),
    ("🎪", "大型活动统筹", "多次参与并协助统筹社团的大型活动，积累丰富的活动组织经验"),
    ("📋", "文件资料整理", "负责部门文件、SRTP资料的收集与整理工作"),
    ("🎉", "文体活动参与", "积极参与社团文体活动，代表社团获太极拳比赛二等奖"),
]

for i, (icon, title, desc) in enumerate(work_items):
    col = i % 2
    row = i // 2
    x = Inches(1.0) + Inches(col * 6.0)
    y = Inches(1.8) + Inches(row * 1.7)

    card = add_glass_card(slide6, x, y, Inches(5.5), Inches(1.4))
    add_textbox(slide6, x + Inches(0.3), y + Inches(0.1), Inches(5.0), Inches(0.4),
                f"{icon}  {title}", font_size=17, color=WHITE, bold=True)
    add_textbox(slide6, x + Inches(0.3), y + Inches(0.6), Inches(4.9), Inches(0.6),
                desc, font_size=12, color=TEXT_GRAY)

# ======================================================================
# SLIDE 7: CLUB CONTRIBUTIONS (社团贡献)
# ======================================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, DARK_BG)

add_textbox(slide7, Inches(1.0), Inches(0.5), Inches(4), Inches(0.6),
            "社团贡献", font_size=36, color=WHITE, bold=True)
add_textbox(slide7, Inches(1.0), Inches(1.0), Inches(4), Inches(0.3),
            "CONTRIBUTIONS", font_size=13, color=ACCENT_BLUE)
add_accent_line(slide7, Inches(1.0), Inches(1.4), Inches(2.0), ACCENT_BLUE, Pt(3))

# Main contribution text
contrib_text = (
    "多次参与并协助统筹社团的大型活动。在担任办公室干事期间，\n"
    "负责日常会议记录、活动物资管理及人员协调。\n"
    "凭借高效的执行力与责任心，顺利完成各项交办任务，\n"
    "并在2026年5月获评社团「优秀干事」。\n\n"
    "此外，也积极参与各项文体活动，曾获太极拳比赛集体二等奖，\n"
    "注重团队协作与个人综合素质的提升。"
)

add_glass_card(slide7, Inches(1.0), Inches(1.9), Inches(7.0), Inches(4.2))
add_multiline_textbox(slide7, Inches(1.5), Inches(2.3), Inches(6.0), Inches(3.5), [
    {'text': '▎核心贡献', 'size': 20, 'color': ACCENT_BLUE, 'bold': True, 'space_after': 12},
    {'text': contrib_text, 'size': 15, 'color': WHITE, 'bold': False, 'space_after': 6},
])

# Right side: key stats
stats = [
    ("优秀干事", "2026年5月获评"),
    ("太极拳二等奖", "2026年4月团体荣誉"),
    ("班级学委", "信安2025-02班"),
]

for i, (num, label) in enumerate(stats):
    y = Inches(2.2) + Inches(i * 1.5)
    add_glass_card(slide7, Inches(8.5), y, Inches(4.0), Inches(1.2))
    add_textbox(slide7, Inches(8.9), y + Inches(0.1), Inches(3.2), Inches(0.5),
                num, font_size=22, color=GOLD, bold=True)
    add_textbox(slide7, Inches(8.9), y + Inches(0.6), Inches(3.2), Inches(0.4),
                label, font_size=13, color=TEXT_GRAY)

# Images
add_image_safe(slide7, "导演.jpg", Inches(8.8), Inches(5.2), width=Inches(1.6))
add_image_safe(slide7, "做饭.jpg", Inches(10.6), Inches(5.2), width=Inches(1.6))

# ======================================================================
# SLIDE 8: REASONS FOR STAYING (留任原因)
# ======================================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, DARK_BG)

add_textbox(slide8, Inches(1.0), Inches(0.5), Inches(4), Inches(0.6),
            "留任原因", font_size=36, color=WHITE, bold=True)
add_textbox(slide8, Inches(1.0), Inches(1.0), Inches(4), Inches(0.3),
            "WHY I STAY", font_size=13, color=ACCENT_BLUE)
add_accent_line(slide8, Inches(1.0), Inches(1.4), Inches(2.0), ACCENT_BLUE, Pt(3))

reasons = [
    ("🏠", "温暖归属，感知温度",
     "双创是温暖的大家庭，欢乐的大集体。在这里结交了许多志同道合的朋友，感受到学长学姐贴心周全的照顾，每次活动都轻松愉快，这份归属感让我渴望留下。"),
    ("📈", "榜样引领，持续学习",
     "双创每位成员都相当优秀，从大家身上能学到很多。学长学姐的经验分享与耐心解惑，为我的成长提供了巨大的帮助。我希望像他们一样，成为后来者的引路人。"),
    ("💪", "责任传承，践行担当",
     "从一名被照顾的干事，到有能力照顾他人。我希望能将双创的温度与精神传承下去，为下一届干事创造同样温暖的成长环境，这是责任，也是使命。"),
]

for i, (icon, title, desc) in enumerate(reasons):
    x = Inches(1.0) + Inches(i * 4.0)
    y = Inches(2.0)
    card = add_glass_card(slide8, x, y, Inches(3.6), Inches(4.8))
    # Top accent bar on card
    add_rect(slide8, x, y, Inches(3.6), Inches(0.06), fill_color=[ACCENT_BLUE, ACCENT_PURPLE, ACCENT_PINK][i])
    add_textbox(slide8, x + Inches(0.25), y + Inches(0.4), Inches(3.1), Inches(0.5),
                f"{icon}  {title}", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide8, x + Inches(0.25), y + Inches(1.2), Inches(3.1), Inches(3.3),
                desc, font_size=13, color=TEXT_GRAY, alignment=PP_ALIGN.LEFT)

# ======================================================================
# SLIDE 9: WORK CAPABILITY (工作能力)
# ======================================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, DARK_BG)

add_textbox(slide9, Inches(1.0), Inches(0.5), Inches(4), Inches(0.6),
            "工作能力 & 竞选优势", font_size=36, color=WHITE, bold=True)
add_textbox(slide9, Inches(1.0), Inches(1.0), Inches(4), Inches(0.3),
            "CAPABILITIES & ADVANTAGES", font_size=13, color=ACCENT_BLUE)
add_accent_line(slide9, Inches(1.0), Inches(1.4), Inches(2.0), ACCENT_BLUE, Pt(3))

# Three columns of capabilities
caps = [
    ("学生工作经验", ACCENT_BLUE, [
        "担任班级学习委员",
        "学风领航中心学业指导部成员",
        "双创办公室干事",
        "获「优秀干事」称号",
        "丰富的组织协调经验",
    ]),
    ("专业能力素养", ACCENT_PURPLE, [
        "信安专业扎实学习",
        "积极参与学科竞赛",
        "广泛涉猎文体活动",
        "英语与综合素质并进",
        "持续学习，不断突破",
    ]),
    ("核心软实力", ACCENT_PINK, [
        "高效执行力与责任心",
        "优秀的跨部门沟通能力",
        "多任务统筹协调能力",
        "团队凝聚与激励能力",
        "积极主动的工作态度",
    ]),
]

for i, (title, color, items) in enumerate(caps):
    x = Inches(1.0) + Inches(i * 4.0)
    y = Inches(1.9)
    add_glass_card(slide9, x, y, Inches(3.6), Inches(5.2))
    add_rect(slide9, x, y, Inches(3.6), Inches(0.06), fill_color=color)
    add_textbox(slide9, x + Inches(0.3), y + Inches(0.3), Inches(3.0), Inches(0.5),
                title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        add_textbox(slide9, x + Inches(0.5), y + Inches(1.1) + Inches(j * 0.7), Inches(2.8), Inches(0.5),
                    f"▸ {item}", font_size=13, color=WHITE)

# ======================================================================
# SLIDE 10: WORK OUTLOOK (工作展望)
# ======================================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10, DARK_BG)

add_textbox(slide10, Inches(1.0), Inches(0.5), Inches(4), Inches(0.6),
            "工作展望", font_size=36, color=WHITE, bold=True)
add_textbox(slide10, Inches(1.0), Inches(1.0), Inches(4), Inches(0.3),
            "FUTURE PLAN", font_size=13, color=ACCENT_BLUE)
add_accent_line(slide10, Inches(1.0), Inches(1.4), Inches(2.0), ACCENT_BLUE, Pt(3))

plans = [
    ("01", "优化工作流程，减负增效", ACCENT_BLUE,
     "致力于减少机械性重复工作，优化办公室日常事务处理流程。让干事们将精力投入到更有价值的活动策划与创新实践中，提升部门整体效能。"),
    ("02", "强化部门协作，精准对接", ACCENT_PURPLE,
     "明确各个环节中办公室与活动部、宣传部等其他部门的对接节点，建立标准化协作流程，确保大型双创活动在物资、场地、人员调配上零失误、高效率。"),
    ("03", "完善培养机制，凝聚团队", ACCENT_PINK,
     "完善「老带新」机制，设计趣味性团建活动与例会形式，增强部门凝聚力。主动关心部员状态，打造有温度、有战斗力的团队文化。"),
]

for i, (num, title, color, desc) in enumerate(plans):
    x = Inches(1.0) + Inches(i * 4.0)
    y = Inches(1.9)
    card = add_glass_card(slide10, x, y, Inches(3.6), Inches(5.0))
    add_rect(slide10, x, y, Inches(3.6), Inches(0.06), fill_color=color)

    # Big number
    add_textbox(slide10, x + Inches(0.3), y + Inches(0.2), Inches(1.5), Inches(1.0),
                num, font_size=52, color=color, bold=True)
    add_textbox(slide10, x + Inches(0.3), y + Inches(1.2), Inches(3.0), Inches(0.5),
                title, font_size=17, color=WHITE, bold=True)
    add_textbox(slide10, x + Inches(0.3), y + Inches(1.8), Inches(3.0), Inches(2.8),
                desc, font_size=12, color=TEXT_GRAY)

# ======================================================================
# SLIDE 11: CLOSING / THANK YOU
# ======================================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide11, DARK_BG)

# Decorative top and bottom bars
add_rect(slide11, Inches(0), Inches(0), Inches(13.333), Inches(0.08), fill_color=ACCENT_BLUE)
add_rect(slide11, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), fill_color=ACCENT_BLUE)

# Decorative circles
c1 = slide11.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-1), Inches(5), Inches(5))
c1.fill.background()
c1.line.color.rgb = ACCENT_BLUE
c1.line.width = Pt(0.5)
c2 = slide11.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), Inches(4), Inches(4), Inches(4))
c2.fill.background()
c2.line.color.rgb = ACCENT_PURPLE
c2.line.width = Pt(0.5)

# Main text
add_textbox(slide11, Inches(2.5), Inches(2.0), Inches(8.5), Inches(1.2),
            "感谢各位的聆听", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_line(slide11, Inches(5.0), Inches(3.2), Inches(3.5), ACCENT_BLUE, Pt(4))

add_textbox(slide11, Inches(2.5), Inches(3.6), Inches(8.5), Inches(0.6),
            "THANKS FOR LISTENING", font_size=18, color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

add_textbox(slide11, Inches(2.5), Inches(4.5), Inches(8.5), Inches(0.5),
            "答辩人：孙  杨", font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide11, Inches(2.5), Inches(5.2), Inches(8.5), Inches(0.4),
            "大学生创新创业俱乐部 · 办公室", font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# === SAVE ===
prs.save(OUTPUT)
print(f"PPT saved to: {os.path.abspath(OUTPUT)}")
print(f"Total slides: {len(prs.slides)}")
